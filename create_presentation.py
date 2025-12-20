"""
Portfolio Presentation Generator
Creates a PowerPoint presentation with the same visual identity as the portfolio website
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Portfolio Colors (Flutter Theme)
FLUTTER_BLUE = RGBColor(1, 181, 248)  # #01b5f8
FLUTTER_DARK_BLUE = RGBColor(2, 86, 155)  # #02569B
DARK_BG = RGBColor(10, 15, 24)  # #0a0f18
LIGHT_TEXT = RGBColor(201, 209, 217)  # #c9d1d9
GRAY_TEXT = RGBColor(139, 148, 158)  # #8b949e

def set_background_color(slide, color):
    """Set slide background to solid color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_with_icon(slide, title_text, top, left, width):
    """Add a styled title with gradient effect"""
    title = slide.shapes.add_textbox(left, top, width, Inches(1))
    text_frame = title.text_frame
    text_frame.text = title_text
    text_frame.word_wrap = True
    
    # Style the title
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(48)
    paragraph.font.bold = True
    paragraph.font.color.rgb = FLUTTER_BLUE
    paragraph.alignment = PP_ALIGN.LEFT
    
    return title

def add_text_box(slide, text, top, left, width, height, font_size=14, color=LIGHT_TEXT, bold=False, align=PP_ALIGN.LEFT):
    """Add a text box with custom styling"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    paragraph.alignment = align
    
    return textbox

def add_rounded_rectangle(slide, top, left, width, height, fill_color, text="", text_color=LIGHT_TEXT, text_size=14):
    """Add a rounded rectangle shape with optional text"""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )
    
    # Style the shape
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = FLUTTER_BLUE
    shape.line.width = Pt(1)
    
    if text:
        text_frame = shape.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(text_size)
        paragraph.font.color.rgb = text_color
        paragraph.alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return shape

def create_portfolio_presentation():
    """Create the complete portfolio presentation"""
    
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Remove default layouts
    blank_layout = prs.slide_layouts[6]  # Blank layout
    
    # ============== Slide 1: Title Slide ==============
    slide1 = prs.slides.add_slide(blank_layout)
    set_background_color(slide1, DARK_BG)
    
    # MA Logo Circle
    logo_circle = add_rounded_rectangle(
        slide1, Inches(1.5), Inches(4), Inches(2), Inches(2),
        FLUTTER_BLUE, "MA", RGBColor(255, 255, 255), 72
    )
    
    # Main Title
    add_text_box(slide1, "Mohamed Adel", Inches(1.2), Inches(0.5), Inches(9), Inches(0.8),
                font_size=60, color=FLUTTER_BLUE, bold=True, align=PP_ALIGN.CENTER)
    
    # Subtitle
    add_text_box(slide1, "Flutter Developer | Cross-Platform Mobile Applications", 
                Inches(2.1), Inches(0.5), Inches(9), Inches(0.5),
                font_size=24, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
    
    # Tagline
    add_text_box(slide1, "Building Beautiful Mobile Experiences", 
                Inches(4.5), Inches(0.5), Inches(9), Inches(0.4),
                font_size=18, color=FLUTTER_BLUE, bold=True, align=PP_ALIGN.CENTER)
    
    # ============== Slide 2: About Me ==============
    slide2 = prs.slides.add_slide(blank_layout)
    set_background_color(slide2, DARK_BG)
    
    # Title
    add_title_with_icon(slide2, "👤 About Me", Inches(0.3), Inches(0.5), Inches(9))
    
    # About Text
    about_text = """A Flutter Developer with +2 years of experience building high-performance 
cross-platform mobile applications.

✓ Proficient in state management (BLoC, Cubit, Provider)
✓ Expert in Firebase integration and RESTful APIs
✓ Applying clean architecture principles
✓ Committed to delivering intuitive user experiences
✓ Strong collaboration and problem-solving skills"""
    
    add_text_box(slide2, about_text, Inches(1.2), Inches(0.5), Inches(5.5), Inches(3.5),
                font_size=16, color=LIGHT_TEXT)
    
    # Stats
    stats_top = Inches(1.5)
    stats_left = Inches(6.5)
    
    # Projects stat
    add_rounded_rectangle(slide2, stats_top, stats_left, Inches(2.8), Inches(1),
                         RGBColor(20, 25, 40), "7\nProjects 🚀", FLUTTER_BLUE, 24)
    
    # Experience stat
    add_rounded_rectangle(slide2, stats_top + Inches(1.3), stats_left, Inches(2.8), Inches(1),
                         RGBColor(20, 25, 40), "1 Year\nExperience 💼", FLUTTER_BLUE, 24)
    
    # Clients stat
    add_rounded_rectangle(slide2, stats_top + Inches(2.6), stats_left, Inches(2.8), Inches(1),
                         RGBColor(20, 25, 40), "2\nHappy Clients 😊", FLUTTER_BLUE, 24)
    
    # ============== Slide 3: Technical Skills ==============
    slide3 = prs.slides.add_slide(blank_layout)
    set_background_color(slide3, DARK_BG)
    
    # Title
    add_title_with_icon(slide3, "💻 Technical Skills", Inches(0.3), Inches(0.5), Inches(9))
    
    # Skills grid
    skills = [
        "Flutter & Dart", "Firebase", "State Management",
        "REST APIs", "Clean Architecture", "Git & GitHub",
        "Responsive Design", "Localization (i18n)", "Testing"
    ]
    
    skill_width = Inches(2.8)
    skill_height = Inches(0.8)
    gap = Inches(0.2)
    start_top = Inches(1.3)
    start_left = Inches(0.5)
    
    for i, skill in enumerate(skills):
        row = i // 3
        col = i % 3
        
        top = start_top + row * (skill_height + gap)
        left = start_left + col * (skill_width + gap)
        
        add_rounded_rectangle(slide3, top, left, skill_width, skill_height,
                             RGBColor(20, 25, 40), skill, LIGHT_TEXT, 14)
    
    # ============== Slide 4: Projects - Part 1 ==============
    slide4 = prs.slides.add_slide(blank_layout)
    set_background_color(slide4, DARK_BG)
    
    # Title
    add_title_with_icon(slide4, "📱 Featured Projects", Inches(0.3), Inches(0.5), Inches(9))
    
    # Project 1: Tasketi
    project1_left = Inches(0.5)
    project1_top = Inches(1.3)
    
    add_rounded_rectangle(slide4, project1_top, project1_left, Inches(4.5), Inches(3.5),
                         RGBColor(20, 25, 40))
    
    add_text_box(slide4, "Tasketi", project1_top + Inches(0.3), project1_left + Inches(0.3),
                Inches(4), Inches(0.4), font_size=24, color=FLUTTER_BLUE, bold=True)
    
    add_text_box(slide4, """A clean and modern Flutter notes app with:
• Offline storage using Hive
• Image support
• Light & Dark themes
• Clean Architecture
• Lottie animations""",
                project1_top + Inches(0.9), project1_left + Inches(0.3),
                Inches(4), Inches(2), font_size=13, color=LIGHT_TEXT)
    
    add_text_box(slide4, "Flutter | Hive | Clean Architecture | Dark Mode",
                project1_top + Inches(3), project1_left + Inches(0.3),
                Inches(4), Inches(0.3), font_size=10, color=FLUTTER_BLUE)
    
    # Project 2: BMI Calculator
    project2_left = Inches(5.2)
    
    add_rounded_rectangle(slide4, project1_top, project2_left, Inches(4.5), Inches(3.5),
                         RGBColor(20, 25, 40))
    
    add_text_box(slide4, "BMI Calculator", project1_top + Inches(0.3), project2_left + Inches(0.3),
                Inches(4), Inches(0.4), font_size=24, color=FLUTTER_BLUE, bold=True)
    
    add_text_box(slide4, """A simple and intuitive BMI calculator:
• Instant BMI computation
• Health category classification
• Clean, modern interface
• Smooth performance
• Material Design""",
                project1_top + Inches(0.9), project2_left + Inches(0.3),
                Inches(4), Inches(2), font_size=13, color=LIGHT_TEXT)
    
    add_text_box(slide4, "Flutter | Dart | Material Design | Health & Fitness",
                project1_top + Inches(3), project2_left + Inches(0.3),
                Inches(4), Inches(0.3), font_size=10, color=FLUTTER_BLUE)
    
    # ============== Slide 5: Projects - Part 2 ==============
    slide5 = prs.slides.add_slide(blank_layout)
    set_background_color(slide5, DARK_BG)
    
    # Title
    add_title_with_icon(slide5, "📚 More Projects", Inches(0.3), Inches(0.5), Inches(9))
    
    # Project 3: Bookia App
    add_rounded_rectangle(slide5, Inches(1.3), Inches(0.5), Inches(4.5), Inches(3.5),
                         RGBColor(20, 25, 40))
    
    add_text_box(slide5, "Bookia App", Inches(1.6), Inches(0.8),
                Inches(4), Inches(0.4), font_size=24, color=FLUTTER_BLUE, bold=True)
    
    add_text_box(slide5, """Complete Flutter bookstore application:
• Live API integration
• Cubit state management
• Login/Signup system
• Shopping cart & wishlist
• SharedPreferences storage
• Clean, modern UI""",
                Inches(2.2), Inches(0.8),
                Inches(4), Inches(2), font_size=13, color=LIGHT_TEXT)
    
    add_text_box(slide5, "Flutter | Cubit | REST API | Authentication | E-commerce",
                Inches(4.3), Inches(0.8),
                Inches(4), Inches(0.3), font_size=10, color=FLUTTER_BLUE)
    
    # More projects coming
    add_rounded_rectangle(slide5, Inches(1.3), Inches(5.2), Inches(4.5), Inches(3.5),
                         RGBColor(20, 25, 40))
    
    add_text_box(slide5, "🚀 More Coming Soon!", Inches(2.5), Inches(5.5),
                Inches(4), Inches(0.5), font_size=28, color=FLUTTER_BLUE, bold=True, align=PP_ALIGN.CENTER)
    
    add_text_box(slide5, "Currently working on exciting new projects!",
                Inches(3.2), Inches(5.5),
                Inches(4), Inches(0.5), font_size=14, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
    
    # ============== Slide 6: Experience ==============
    slide6 = prs.slides.add_slide(blank_layout)
    set_background_color(slide6, DARK_BG)
    
    # Title
    add_title_with_icon(slide6, "💼 Experience", Inches(0.3), Inches(0.5), Inches(9))
    
    # Experience 1: NTI
    exp1_top = Inches(1.3)
    add_rounded_rectangle(slide6, exp1_top, Inches(0.5), Inches(9), Inches(1.3),
                         RGBColor(20, 25, 40))
    
    add_text_box(slide6, "National Telecommunication Institute (NTI)",
                exp1_top + Inches(0.15), Inches(0.8), Inches(8), Inches(0.4),
                font_size=18, color=FLUTTER_BLUE, bold=True)
    
    add_text_box(slide6, "Flutter Trainee • Jul 2025 - Aug 2025 • Cairo, Egypt",
                exp1_top + Inches(0.45), Inches(0.8), Inches(8), Inches(0.3),
                font_size=12, color=GRAY_TEXT)
    
    add_text_box(slide6, "• Intensive Flutter training • API integration • BLoC & Cubit • Clean Architecture",
                exp1_top + Inches(0.8), Inches(0.8), Inches(8), Inches(0.4),
                font_size=11, color=LIGHT_TEXT)
    
    # Experience 2: GDSC
    exp2_top = Inches(2.8)
    add_rounded_rectangle(slide6, exp2_top, Inches(0.5), Inches(9), Inches(1.3),
                         RGBColor(20, 25, 40))
    
    add_text_box(slide6, "Google Developer Student Clubs (GDSC)",
                exp2_top + Inches(0.15), Inches(0.8), Inches(8), Inches(0.4),
                font_size=18, color=FLUTTER_BLUE, bold=True)
    
    add_text_box(slide6, "Flutter Trainee • Jan 2024 - Feb 2024 • Benha University",
                exp2_top + Inches(0.45), Inches(0.8), Inches(8), Inches(0.3),
                font_size=12, color=GRAY_TEXT)
    
    add_text_box(slide6, "• 36-hour bootcamp • UI development • OOP principles • Calculator App project",
                exp2_top + Inches(0.8), Inches(0.8), Inches(8), Inches(0.4),
                font_size=11, color=LIGHT_TEXT)
    
    # Experience 3: D.Pear
    exp3_top = Inches(4.3)
    add_rounded_rectangle(slide6, exp3_top, Inches(0.5), Inches(9), Inches(1.3),
                         RGBColor(20, 25, 40))
    
    add_text_box(slide6, "D.Pear",
                exp3_top + Inches(0.15), Inches(0.8), Inches(8), Inches(0.4),
                font_size=18, color=FLUTTER_BLUE, bold=True)
    
    add_text_box(slide6, "Flutter Summer Trainee • Summer 2024 (6 Weeks) • Remote",
                exp3_top + Inches(0.45), Inches(0.8), Inches(8), Inches(0.3),
                font_size=12, color=GRAY_TEXT)
    
    add_text_box(slide6, "• Multiple projects • API & Firebase integration • Hive database • Clean Architecture",
                exp3_top + Inches(0.8), Inches(0.8), Inches(8), Inches(0.4),
                font_size=11, color=LIGHT_TEXT)
    
    # ============== Slide 7: Contact ==============
    slide7 = prs.slides.add_slide(blank_layout)
    set_background_color(slide7, DARK_BG)
    
    # Title
    add_title_with_icon(slide7, "📧 Get In Touch", Inches(0.3), Inches(0.5), Inches(9))
    
    # Contact cards
    contact_top = Inches(1.5)
    contact_left = Inches(1)
    card_width = Inches(4)
    card_height = Inches(0.9)
    gap = Inches(0.25)
    
    # Location
    add_rounded_rectangle(slide7, contact_top, contact_left, card_width, card_height,
                         RGBColor(20, 25, 40))
    add_text_box(slide7, "📍 Location", contact_top + Inches(0.15), contact_left + Inches(0.3),
                Inches(3.5), Inches(0.3), font_size=14, color=FLUTTER_BLUE, bold=True)
    add_text_box(slide7, "Qaliobia, Benha, Egypt", contact_top + Inches(0.5), contact_left + Inches(0.3),
                Inches(3.5), Inches(0.3), font_size=12, color=LIGHT_TEXT)
    
    # Phone
    add_rounded_rectangle(slide7, contact_top, contact_left + card_width + gap, card_width, card_height,
                         RGBColor(20, 25, 40))
    add_text_box(slide7, "📱 Phone", contact_top + Inches(0.15), contact_left + card_width + gap + Inches(0.3),
                Inches(3.5), Inches(0.3), font_size=14, color=FLUTTER_BLUE, bold=True)
    add_text_box(slide7, "+20 106 412 0753", contact_top + Inches(0.5), contact_left + card_width + gap + Inches(0.3),
                Inches(3.5), Inches(0.3), font_size=12, color=LIGHT_TEXT)
    
    # Email
    contact_top2 = contact_top + card_height + gap
    add_rounded_rectangle(slide7, contact_top2, contact_left, card_width, card_height,
                         RGBColor(20, 25, 40))
    add_text_box(slide7, "✉️ Email", contact_top2 + Inches(0.15), contact_left + Inches(0.3),
                Inches(3.5), Inches(0.3), font_size=14, color=FLUTTER_BLUE, bold=True)
    add_text_box(slide7, "mahmedadel973@gmail.com", contact_top2 + Inches(0.5), contact_left + Inches(0.3),
                Inches(3.5), Inches(0.3), font_size=11, color=LIGHT_TEXT)
    
    # LinkedIn
    add_rounded_rectangle(slide7, contact_top2, contact_left + card_width + gap, card_width, card_height,
                         RGBColor(20, 25, 40))
    add_text_box(slide7, "💼 LinkedIn", contact_top2 + Inches(0.15), contact_left + card_width + gap + Inches(0.3),
                Inches(3.5), Inches(0.3), font_size=14, color=FLUTTER_BLUE, bold=True)
    add_text_box(slide7, "Mohamed Adel", contact_top2 + Inches(0.5), contact_left + card_width + gap + Inches(0.3),
                Inches(3.5), Inches(0.3), font_size=12, color=LIGHT_TEXT)
    
    # Social Links section
    add_text_box(slide7, "Connect with me:", Inches(3.8), Inches(1), Inches(8), Inches(0.3),
                font_size=16, color=FLUTTER_BLUE, bold=True, align=PP_ALIGN.CENTER)
    
    social_top = Inches(4.2)
    social_left = Inches(2)
    social_gap = Inches(1.5)
    
    socials = [
        ("🔗 GitHub", "github.com/MohamedAdel743"),
        ("💼 LinkedIn", "linkedin.com/in/mohamed-adel"),
        ("📘 Facebook", "facebook.com/modola123"),
        ("💬 WhatsApp", "+20 106 412 0753")
    ]
    
    for i, (label, link) in enumerate(socials):
        add_rounded_rectangle(slide7, social_top, social_left + i * social_gap, Inches(1.3), Inches(0.7),
                             RGBColor(20, 25, 40))
        add_text_box(slide7, f"{label}\n{link}", social_top + Inches(0.1), social_left + i * social_gap + Inches(0.1),
                    Inches(1.1), Inches(0.5), font_size=8, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)
    
    # Save presentation
    output_path = os.path.join(os.path.dirname(__file__), "Mohamed_Adel_Portfolio_Presentation.pptx")
    prs.save(output_path)
    print(f"✅ Presentation created successfully!")
    print(f"📁 Saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    try:
        create_portfolio_presentation()
    except Exception as e:
        print(f"❌ Error creating presentation: {str(e)}")
        raise
